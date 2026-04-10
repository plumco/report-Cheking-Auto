"""
Huliot India — Site Visit Report Formatter + AI Comment Improver
================================================================
• Upload team PPTX
• Auto-fix: font sizes, time/date format, banner text, member labels, image sizes
• AI rewrites observation comments to match Huliot professional standard
• You review & edit comments before downloading
• Download corrected, perfectly formatted PPTX

Run: streamlit run app.py
"""

import streamlit as st
import zipfile, re, io, time, json, copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
import requests

# ─────────────────────────────────────────────────────────────────────────────
#  CONSTANTS — HULIOT STANDARD FORMAT
# ─────────────────────────────────────────────────────────────────────────────

# EMU sizes
EMU = 914400                   # 1 inch in EMU

SLIDE_W   = int(13.33 * EMU)  # 13.33"
SLIDE_H   = int(7.50  * EMU)  # 7.50"

# Standard font sizes (pt) per role
FONT_COVER_TITLE  = 32   # "Huliot India:- ProjectName"
FONT_COVER_BANNER = 28   # "Site Visit" green banner
FONT_COVER_BODY   = 20   # bullets under banner
FONT_EXPLAIN_HEAD = 32   # "Explain - Installation instruction"
FONT_OBS_HEAD     = 24   # Observation slide heading
FONT_CONTENT      = 18   # Body text on explain/checklist slides
FONT_THANKYOU     = 16   # Thank you slide

# Standard image sizes (inches) for observation slide photos
OBS_PHOTO_W   = 4.00   # site photo width
OBS_PHOTO_H   = 4.00   # site photo height
OBS_DRAW_W    = 3.50   # drawing snapshot width
OBS_DRAW_H    = 3.50   # drawing snapshot height

# ─────────────────────────────────────────────────────────────────────────────
#  STRICT-OOXML FIX (purl.oclc.org → schemas.openxmlformats.org)
# ─────────────────────────────────────────────────────────────────────────────

NS_MAP = {
    "http://purl.oclc.org/ooxml/drawingml/main":
        "http://schemas.openxmlformats.org/drawingml/2006/main",
    "http://purl.oclc.org/ooxml/officeDocument/relationships":
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "http://purl.oclc.org/ooxml/presentationml/main":
        "http://schemas.openxmlformats.org/presentationml/2006/main",
    "http://purl.oclc.org/ooxml/officeDocument/relationships/officeDocument":
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument",
    "http://purl.oclc.org/ooxml/officeDocument/relationships/extendedProperties":
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties",
}

def fix_strict_ooxml(pptx_bytes: bytes) -> bytes:
    buf, out = io.BytesIO(pptx_bytes), io.BytesIO()
    with zipfile.ZipFile(buf, "r") as zin:
        with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith(".xml") or item.filename.endswith(".rels"):
                    try:
                        text = data.decode("utf-8")
                        for old, new in NS_MAP.items():
                            text = text.replace(old, new)
                        data = text.encode("utf-8")
                    except Exception:
                        pass
                zout.writestr(item, data)
    out.seek(0)
    return out.read()

# ─────────────────────────────────────────────────────────────────────────────
#  HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return " ".join("".join(r.text for r in p.runs) for p in shape.text_frame.paragraphs).strip()


def set_para_text(para, new_text: str):
    """Replace paragraph text into first run, clear others."""
    if not para.runs:
        return
    para.runs[0].text = new_text
    for r in para.runs[1:]:
        r.text = ""


def fix_para_pattern(para, pattern: str, repl, flags: int = 0) -> tuple:
    full = "".join(r.text for r in para.runs)
    m = re.search(pattern, full, flags)
    if not m:
        return False, "", ""
    old = m.group(0)
    if callable(repl):
        new_full = re.sub(pattern, repl, full, flags=flags)
    else:
        new_full = re.sub(pattern, repl, full, flags=flags)
    if para.runs:
        para.runs[0].text = new_full
        for r in para.runs[1:]:
            r.text = ""
    return True, old, new_full


def set_font_size(shape, target_pt: float):
    """Set all runs in a shape to target font size."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:  # only if explicitly set
                run.font.size = Pt(target_pt)


def clamp_image_size(shape, max_w_in: float, max_h_in: float):
    """Resize image to fit within max dimensions, preserving aspect ratio."""
    if shape.shape_type != 13:
        return False
    w, h = shape.width, shape.height
    max_w = int(max_w_in * EMU)
    max_h = int(max_h_in * EMU)
    if w <= max_w and h <= max_h:
        return False
    ratio = min(max_w / w, max_h / h)
    shape.width  = int(w * ratio)
    shape.height = int(h * ratio)
    return True

# ─────────────────────────────────────────────────────────────────────────────
#  FORMAT FIXER — Applies all standard formatting rules
# ─────────────────────────────────────────────────────────────────────────────

def apply_format_fixes(prs: Presentation) -> list:
    """
    Apply all Huliot standard format fixes to a Presentation object.
    Returns list of fix descriptions.
    """
    fixes = []
    N = len(prs.slides)

    # ── SLIDE 1: COVER ──────────────────────────────────────────────────────
    if N >= 1:
        slide = prs.slides[0]
        W, H = prs.slide_width, prs.slide_height

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full = shape_text(shape)

            # Green banner detection → must say "Site Visit"
            is_banner = (
                shape.width > W * 0.3 and
                shape.height < H * 0.15 and
                H * 0.25 < (shape.top + shape.height / 2) < H * 0.65 and
                len(shape.text_frame.paragraphs) <= 2 and
                full and "Huliot" not in full and len(full) < 80
            )
            if is_banner and "Site Visit" not in full:
                first = True
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = "Site Visit" if first else ""
                        if first:
                            run.font.size = Pt(FONT_COVER_BANNER)
                        first = False
                fixes.append(("Slide 1", "Banner text", f'"{full}" → "Site Visit"'))

            # Ensure banner font size
            if "Site Visit" in full and is_banner:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size != Pt(FONT_COVER_BANNER):
                            run.font.size = Pt(FONT_COVER_BANNER)

            # Time format fix
            if re.search(r"\d+[.:]\d+\s*[Tt][oO]\s*\d+[.:]\d+", full):
                for para in shape.text_frame.paragraphs:
                    ch, old, _ = fix_para_pattern(
                        para,
                        r"(\d{1,2})[.:](\d{2})\s*[Tt][oO]\s*(\d{1,2})[.:](\d{2})",
                        lambda m: f"{m.group(1)}:{m.group(2)}am to {m.group(3)}:{m.group(4)}pm"
                    )
                    if ch:
                        fixes.append(("Slide 1", "Time format", f'"{old}" → correct am/pm format'))

            # Date spacing
            if re.search(r"\d+\s+\.\d+", full):
                for para in shape.text_frame.paragraphs:
                    fix_para_pattern(para, r"(\d+)\s+\.(\d+)", r"\1.\2")
                fixes.append(("Slide 1", "Date spacing", "Extra space before dot removed"))

            # "Client" → "Plumbers"
            if re.search(r"\bClient\s*:", full, re.IGNORECASE):
                for para in shape.text_frame.paragraphs:
                    fix_para_pattern(para, r"\bClient\s*:", "Plumbers :-", flags=re.IGNORECASE)
                fixes.append(("Slide 1", "Members label", '"Client :" → "Plumbers :-"'))

            # Font size enforcement on cover
            for para in shape.text_frame.paragraphs:
                full_para = "".join(r.text for r in para.runs)
                for run in para.runs:
                    if run.font.size:
                        # Title block: 32pt
                        if any(kw in full_para for kw in ["Huliot India:-", "Date:-", "Time:-"]):
                            if abs(run.font.size - Pt(FONT_COVER_TITLE)) > Pt(2):
                                run.font.size = Pt(FONT_COVER_TITLE)
                        # Body bullets: 20pt
                        elif any(kw in full_para for kw in ["Site Name", "Location", "Members", "Contractor", "Plumbers", "Huliot India –"]):
                            if abs(run.font.size - Pt(FONT_COVER_BODY)) > Pt(2):
                                run.font.size = Pt(FONT_COVER_BODY)

    # ── SLIDES 2+ ───────────────────────────────────────────────────────────
    for idx in range(1, N):
        snum  = idx + 1
        slide = prs.slides[idx]

        # Image size standardization
        pics = [s for s in slide.shapes if s.shape_type == 13]
        for pic in pics:
            w_in = pic.width  / EMU
            h_in = pic.height / EMU
            # Oversized: anything wider than 5.5" or taller than 5"
            if w_in > 5.5 or h_in > 5.0:
                old_w, old_h = w_in, h_in
                clamped = clamp_image_size(pic, 5.0, 4.5)
                if clamped:
                    fixes.append((f"Slide {snum}", "Image size",
                                  f"Photo resized: {old_w:.1f}\" x {old_h:.1f}\" → max 5\"x4.5\""))

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full = shape_text(shape)

            # Spelling: PRL ines → PRL lines
            if re.search(r"\bPRL\s+ines\b", full, re.IGNORECASE):
                for para in shape.text_frame.paragraphs:
                    fix_para_pattern(para, r"(PRL\s+)ines\b", r"\1lines", flags=re.IGNORECASE)
                fixes.append((f"Slide {snum}", "Spelling", '"PRL ines" → "PRL lines"'))

            # Font enforcement on Explain slides (headings)
            if any(kw in full for kw in ["Explain", "Drainage Checklist"]):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            target = FONT_OBS_HEAD if snum == 2 else FONT_EXPLAIN_HEAD
                            if run.font.size < Pt(target - 4) or run.font.size > Pt(target + 6):
                                run.font.size = Pt(target)

            # Font enforcement on content text (18pt)
            if snum not in [1, 11] and "Explain" not in full and "Drainage Checklist" not in full:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.text.strip() and len(run.text) > 5:
                            if run.font.size > Pt(22) and "u@" not in run.text:
                                run.font.size = Pt(FONT_CONTENT)

        # Slide 11 — Thank You font
        if snum == 11:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size > Pt(20):
                                run.font.size = Pt(FONT_THANKYOU)

    return fixes

# ─────────────────────────────────────────────────────────────────────────────
#  EXTRACT OBSERVATION COMMENTS from slide 2 (the real site observations)
# ─────────────────────────────────────────────────────────────────────────────

def extract_obs_comments(prs: Presentation) -> list:
    """
    Extract text blocks from observation slides (slide 2 onwards, until standard slides).
    Returns list of {slide, shape_idx, text}.
    """
    obs = []
    N = len(prs.slides)
    # Observation slides are those that aren't "Explain" or "Thank You" or checklist
    SKIP_KEYWORDS = ["Explain", "Drainage Checklist", "Thank", "PRL Option", "Site Visit"]
    for idx in range(1, min(N - 1, 5)):  # Check slides 2-5
        slide = prs.slides[idx]
        for shi, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            full = shape_text(shape).strip()
            if not full or len(full) < 20:
                continue
            if any(kw in full for kw in SKIP_KEYWORDS):
                continue
            if "u@" in full:
                continue
            obs.append({"slide": idx + 1, "shape_idx": shi, "text": full})
    return obs

# ─────────────────────────────────────────────────────────────────────────────
#  AI COMMENT IMPROVER — Calls Claude API
# ─────────────────────────────────────────────────────────────────────────────

SYSTEM_PROMPT = """You are a senior Huliot Technical Manager (India) with 15 years experience.
You review site visit observation text written by junior field engineers and rewrite it 
to professional standard. 

Rules:
- Keep the same facts — do NOT invent or remove observations
- Fix grammar, spelling, and punctuation
- Use formal English suitable for client report
- Start with ✅ (positive) OR ❌ (non-compliance) where appropriate
- Keep correction instructions clear and actionable
- Maximum 3 sentences per observation
- Do NOT add bullet points — single flowing paragraph per observation
- Reference Huliot standards where relevant (clamp spacing, 45° bend, PRL clearance, etc.)

Return ONLY the improved text. No preamble, no explanation."""


def ai_improve_comment(raw_text: str) -> str:
    """Send raw comment to Claude API and return improved version."""
    try:
        resp = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"Content-Type": "application/json"},
            json={
                "model": "claude-sonnet-4-20250514",
                "max_tokens": 400,
                "system": SYSTEM_PROMPT,
                "messages": [{"role": "user", "content": f"Improve this site observation comment:\n\n{raw_text}"}]
            },
            timeout=30
        )
        data = resp.json()
        for block in data.get("content", []):
            if block.get("type") == "text":
                return block["text"].strip()
    except Exception as e:
        return raw_text  # fallback: return original
    return raw_text


def ai_improve_all_comments(obs_list: list) -> list:
    """Improve all observation comments via AI."""
    improved = []
    for obs in obs_list:
        ai_text = ai_improve_comment(obs["text"])
        improved.append({**obs, "ai_text": ai_text})
    return improved

# ─────────────────────────────────────────────────────────────────────────────
#  APPLY FINAL COMMENTS back into PPTX
# ─────────────────────────────────────────────────────────────────────────────

def apply_comments(prs: Presentation, final_comments: list):
    """Write final (user-approved) comments back into the PPTX."""
    for item in final_comments:
        slide = prs.slides[item["slide"] - 1]
        shapes = list(slide.shapes)
        if item["shape_idx"] < len(shapes):
            shape = shapes[item["shape_idx"]]
            if shape.has_text_frame and shape.text_frame.paragraphs:
                # Write into first paragraph first run
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    para.runs[0].text = item["final_text"]
                    for r in para.runs[1:]:
                        r.text = ""
                    # Clear extra paragraphs if needed (keep 1)
                    for extra_para in shape.text_frame.paragraphs[1:]:
                        for r in extra_para.runs:
                            r.text = ""

# ─────────────────────────────────────────────────────────────────────────────
#  STREAMLIT APP
# ─────────────────────────────────────────────────────────────────────────────

def main():
    st.set_page_config(
        page_title="Huliot Report Formatter",
        page_icon="🏗️",
        layout="wide",
    )

    # Header
    st.markdown("""
    <div style="background:linear-gradient(135deg,#1B5E20 0%,#2E7D32 60%,#388E3C 100%);
                padding:24px 32px;border-radius:14px;margin-bottom:24px;
                box-shadow:0 4px 16px rgba(0,0,0,0.2);">
        <div style="color:#A5D6A7;font-size:11px;font-weight:700;letter-spacing:3px;
                    text-transform:uppercase;margin-bottom:6px;">
            Huliot Pipes & Fittings Pvt. Ltd. — West Zone
        </div>
        <div style="color:white;font-size:26px;font-weight:800;line-height:1.2;">
            🏗️ Site Visit Report — Smart Formatter & AI Comment Fixer
        </div>
        <div style="color:#C8E6C9;font-size:13px;margin-top:8px;">
            Upload team PPT → Auto-fix format & font → AI improves comments → You review → Download perfect report
        </div>
    </div>
    """, unsafe_allow_html=True)

    # ── Sidebar: what this fixes ──────────────────────────────────────────────
    with st.sidebar:
        st.markdown("### 📋 Standard Format Rules")
        st.markdown("""
        **Cover Slide**
        - ✅ Banner = "Site Visit"
        - ✅ Time: `10:30am to 01:30pm`
        - ✅ Title font: **32pt**
        - ✅ Body bullets: **20pt**
        - ✅ Label: "Plumbers :-" (not "Client")

        **Observation Slides**
        - ✅ Heading font: **24pt**
        - ✅ Body text: **18pt**
        - ✅ Max photo width: **5 inches**
        - ✅ Max photo height: **4.5 inches**
        - ✅ Spelling errors auto-fixed

        **Explain Slides (3–8)**
        - ✅ Heading font: **32pt**
        - ✅ Content text: **18pt**

        **Checklist Slide 9**
        - ✅ Heading font: **32pt**

        **AI Comment Improvement**
        - ✅ Professional English
        - ✅ Huliot standards referenced
        - ✅ Clear correction instructions
        - ✅ You review & edit before download
        """)

    # ── Upload ────────────────────────────────────────────────────────────────
    st.markdown("### 📎 Step 1 — Upload Team Report")
    uploaded = st.file_uploader(
        "Drop .pptx file here",
        type=["pptx"],
        label_visibility="collapsed"
    )

    if not uploaded:
        st.info("👆 Upload a **.pptx** report from your team to begin.")
        return

    st.success(f"✅ **{uploaded.name}** loaded — {uploaded.size/1024:.1f} KB")
    st.divider()

    # ── Process button ────────────────────────────────────────────────────────
    st.markdown("### ⚙️ Step 2 — Format Fix + AI Comment Improvement")
    run_btn = st.button(
        "▶ Fix Format & Improve Comments (AI)",
        type="primary",
        use_container_width=True
    )

    if run_btn or st.session_state.get("processed"):

        if run_btn:
            # ── Do the processing ─────────────────────────────────────────────
            pptx_bytes = uploaded.read()
            bar    = st.progress(0)
            status = st.empty()

            # Step A: Fix namespace
            bar.progress(10, "Fixing file format...")
            safe_bytes = fix_strict_ooxml(pptx_bytes)

            # Step B: Load & apply format fixes
            bar.progress(30, "Applying standard format (fonts, sizes, labels)...")
            prs = Presentation(io.BytesIO(safe_bytes))
            format_fixes = apply_format_fixes(prs)

            # Step C: Extract observation comments
            bar.progress(50, "Extracting observation comments...")
            obs_list = extract_obs_comments(prs)

            # Step D: AI improve comments
            bar.progress(65, f"AI improving {len(obs_list)} observation comment(s)...")
            improved = ai_improve_all_comments(obs_list)

            # Step E: Save prs to bytes (pre-comment state)
            bar.progress(85, "Saving formatted file...")
            pre_out = io.BytesIO()
            prs.save(pre_out)
            pre_out.seek(0)
            pre_bytes = pre_out.read()

            bar.progress(100, "Done!")
            time.sleep(0.4)
            bar.empty()
            status.empty()

            # Store in session
            st.session_state["processed"]    = True
            st.session_state["pre_bytes"]    = pre_bytes
            st.session_state["format_fixes"] = format_fixes
            st.session_state["improved"]     = improved
            st.session_state["filename"]     = uploaded.name

        # ── Show format fix results ────────────────────────────────────────────
        st.divider()
        st.markdown("### ✅ Format Fixes Applied")

        ff = st.session_state.get("format_fixes", [])
        if ff:
            cols = st.columns(3)
            for i, (slide, field, desc) in enumerate(ff):
                with cols[i % 3]:
                    st.success(f"**{slide}** — {field}\n\n{desc}")
        else:
            st.info("No format issues found — file was already correct.")

        st.divider()

        # ── AI Comment Review ──────────────────────────────────────────────────
        improved = st.session_state.get("improved", [])
        st.markdown("### 🤖 Step 3 — Review & Edit AI-Improved Comments")

        if not improved:
            st.info("No observation comments found to improve on slides 2–5.")
        else:
            st.caption(
                "👇 AI has rewritten each comment to professional standard. "
                "**You can edit any text below before downloading.**"
            )

            final_comments = []
            for i, item in enumerate(improved):
                with st.expander(
                    f"📝 Slide {item['slide']} — Comment {i+1}",
                    expanded=True
                ):
                    col_orig, col_ai = st.columns(2)
                    with col_orig:
                        st.markdown("**🔴 Original (team wrote):**")
                        st.text_area(
                            "original",
                            value=item["text"],
                            height=120,
                            disabled=True,
                            label_visibility="collapsed",
                            key=f"orig_{i}"
                        )
                    with col_ai:
                        st.markdown("**✅ AI Improved (editable):**")
                        edited = st.text_area(
                            "edited",
                            value=item.get("ai_text", item["text"]),
                            height=120,
                            label_visibility="collapsed",
                            key=f"edit_{i}"
                        )
                    final_comments.append({
                        "slide":      item["slide"],
                        "shape_idx":  item["shape_idx"],
                        "final_text": edited,
                    })

            st.session_state["final_comments"] = final_comments

        st.divider()

        # ── Download ───────────────────────────────────────────────────────────
        st.markdown("### ⬇️ Step 4 — Download Final Corrected Report")

        if st.button("📦 Generate Final File", type="secondary", use_container_width=True):
            # Re-load, re-fix, apply final comments
            pre_bytes = st.session_state["pre_bytes"]
            prs2 = Presentation(io.BytesIO(pre_bytes))
            final_comments = st.session_state.get("final_comments", [])
            if final_comments:
                apply_comments(prs2, final_comments)
            final_out = io.BytesIO()
            prs2.save(final_out)
            final_out.seek(0)
            final_bytes = final_out.read()

            orig_name = st.session_state.get("filename", "report.pptx")
            out_name = f"FORMATTED_{Path(orig_name).stem}.pptx"

            st.download_button(
                label=f"⬇️  Download: {out_name}",
                data=final_bytes,
                file_name=out_name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                type="primary",
            )
            st.caption(
                "✅ File ready. Open in PowerPoint → review once → send to client."
            )

        else:
            # Quick download without re-generating (uses pre-comment bytes)
            st.caption("👆 Click above to apply your comment edits and generate the final file.")


if __name__ == "__main__":
    main()
